import re
from pptx import Presentation


def clean_text(text):
    if not text:
        return ""
    return "".join(ch for ch in str(text) if ch.isprintable())


def extract_azure_bug(text):
    match = re.search(r"\b(\d{5,})\b", text or "")
    return match.group(1) if match else ""


def extract_slide1_metadata(ppt_path):
    prs = Presentation(ppt_path)

    if not prs.slides:
        return {}

    slide = prs.slides[0]
    texts = []

    for shape in slide.shapes:
        if hasattr(shape, "text"):
            txt = clean_text(shape.text.strip())
            if txt:
                texts.append(txt)

    full_text = " ".join(texts)

    incident_match = re.search(r"(INC\d+)", full_text)
    date_match = re.search(r"\d{1,2}-[A-Za-z]{3}-\d{4}", full_text)

    return {
        "incident": incident_match.group(1) if incident_match else "",
        "description": full_text,
        "created_date": date_match.group(0) if date_match else "",
        "azure_bug": extract_azure_bug(full_text)
    }
