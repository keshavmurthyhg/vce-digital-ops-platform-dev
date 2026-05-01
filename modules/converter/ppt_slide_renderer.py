import os
import tempfile
import uuid
from pptx import Presentation


def render_ppt_slides_to_images(ppt_path):
    """
    Extract largest image from each slide.
    No Aspose dependency.
    Keeps app stable on Streamlit Cloud.
    """

    prs = Presentation(ppt_path)
    output_dir = tempfile.mkdtemp()

    slide_images = []

    for slide_index, slide in enumerate(prs.slides):

        largest_image = None
        largest_size = 0

        for shape in slide.shapes:
            try:
                if shape.shape_type == 13:  # Picture
                    image_bytes = shape.image.blob
                    image_size = len(image_bytes)

                    if image_size > largest_size:
                        largest_size = image_size
                        largest_image = image_bytes

            except Exception:
                continue

        if largest_image:
            img_path = os.path.join(
                output_dir,
                f"slide_{slide_index+1}_{uuid.uuid4().hex}.png"
            )

            with open(img_path, "wb") as f:
                f.write(largest_image)

            slide_images.append(img_path)

    return slide_images
