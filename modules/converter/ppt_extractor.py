from pptx import Presentation
import os
import uuid


def extract_ppt_content(ppt_path, tmp_dir):
    prs = Presentation(ppt_path)

    slides_data = []

    for i, slide in enumerate(prs.slides):

        # Skip slide 1 (already used in header)
        if i == 0:
            continue

        slide_info = {
            "title": f"Slide {i+1}",
            "texts": [],
            "images": []
        }

        for shape in slide.shapes:

            # TEXT
            if hasattr(shape, "text") and shape.text.strip():
                slide_info["texts"].append(shape.text.strip())

            # IMAGE
            try:
                if shape.shape_type == 13:  # picture
                    img_path = os.path.join(tmp_dir, f"{uuid.uuid4().hex}.png")

                    with open(img_path, "wb") as f:
                        f.write(shape.image.blob)

                    slide_info["images"].append(img_path)

            except Exception:
                continue

        slides_data.append(slide_info)

    return slides_data
